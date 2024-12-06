import fitz  # PyMuPDF
import os
from groq import Groq
import base64
import logging
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tempfile
from pathlib import Path



def api_setup():
    """
    Sets up the Groq API client using the GROQ_API_KEY environment variable.

    Returns:
        client: An initialized Groq client object.

    Raises:
        ValueError: If the GROQ_API_KEY environment variable is not set.
    """
    GROQ_API_KEY = os.getenv('GROQ_API_KEY')
    if GROQ_API_KEY is None:
        raise ValueError("The GROQ_API_KEY environment variable is not set.")
    # Initialize the Groq client
    client = Groq(api_key=GROQ_API_KEY)
    return client



def encode_image(image_stream):
    """
    Encodes an image stream to a base64 string.

    Args:
        image_stream (io.BytesIO): A binary stream of the image to be encoded.

    Returns:
        str: The base64 encoded string of the image.
    """
    return base64.b64encode(image_stream.read()).decode('utf-8')



def save_text(text_to_save, file_name="full_text.txt"):
    """
    Saves the given text to a file, replacing specific placeholder characters with their corresponding German umlaut characters.

    Args:
        text_to_save (str): The text to be saved, with placeholders for umlaut characters.
        file_name (str, optional): The name of the file to save the text to. Defaults to "full_text.txt".

    Returns:
        None
    """
    text_to_save = text_to_save.replace("¨a", "ä").replace("¨o", "ö").replace("¨u", "ü").replace("¨A", "Ä").replace("¨O", "Ö").replace("¨U", "Ü")
    print("save function running")
    full_text_path = Path(__file__).parent / file_name
    if not full_text_path.exists():
        full_text_path.touch()
    with open(full_text_path, "a", encoding="UTF-8") as file:
        file.write(text_to_save)




def describe_image(image_stream, prompt):
    """
    Generates a description of an image based on a given prompt using an AI model.

    Args:
        image_stream (io.BytesIO): The image data stream to be described.
        prompt (str): The text prompt to guide the description.

    Returns:
        str: The generated description of the image, or None if an error occurs.

    Raises:
        Exception: If there is an error during the API request.
    """
    client = api_setup()
    base64_image = encode_image(image_stream)
    try:
        chat_completion = client.chat.completions.create(
                        messages=[
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": prompt},
                                {"type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{base64_image}",
                            },},],}],
                    model="llama-3.2-90b-vision-preview",)
        return chat_completion.choices[0].message.content
    except Exception as e:
        logging.error(f"Error during summary API request: {e}")
        return None



def summarize(long_text, prompt):
    """
    Summarizes a given long text using a specified prompt.

    Args:
        long_text (str): The text to be summarized.
        prompt (str): The prompt to guide the summarization.

    Returns:
        str: The summarized text if the API request is successful.
        None: If there is an error during the API request.

    Raises:
        Exception: If there is an error during the summary API request.
    """
    client = api_setup()
    try:
        chat_completion = client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": prompt,
                    },
                    {
                        "role": "user",
                        "content": long_text,
                    },
                ],
                model="llama3-8b-8192",
            )
        response = chat_completion.choices[0].message.content
        return response
    except Exception as e:
        logging.error(f"Error during summary API request: {e}")
        return None



def describe_pdf(pdf, short_response=False, max_length=3000, store_content=False):
    """
    Extracts text and images from a PDF file, summarizes the content, and optionally stores the extracted text.

    Args:
        pdf (file-like object): The PDF file to be described.
        short_response (bool, optional): If True, the summary will be shortened if it exceeds max_length. Defaults to False.
        max_length (int, optional): The maximum length of the summary. Defaults to 3000.
        store_content (bool, optional): If True, the extracted text will be saved to a file. Defaults to False.

    Returns:
        str: The summarized content of the PDF.
    """
    pdf_stream = io.BytesIO(pdf.read())
    print(pdf.name)
    image_prompt = "Das folgende Bild ist Teil einer Presentation, fasse zentrale Inhalte des Bildes in maximal 100 Worten zusammen, ohne Information hinzuzufügen"
    summary_prompt = "du fasst die wichtigsten informationen eines Textes auf deutsch in ganzen sätzen zusammen. Füge keine Informationen hinzu. gebe nur die Zusammenfassung zurück. Füge keine Textformatierung hinzu."
    pdf_document = fitz.open(stream=pdf_stream, filetype="pdf")
    response = ""
    for page_num in range(pdf_document.page_count):
        text = ""
        page = pdf_document[page_num]
        text += page.get_text().encode("utf-8").decode("utf-8")
        image_list = page.get_images(full=True)
        for img in image_list:
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_stream = io.BytesIO(image_bytes)
            text += describe_image(image_stream, image_prompt)
        response += summarize(text, summary_prompt)
        if store_content:
            file_name = str(os.path.splitext(pdf.name)[0])+".txt"
            save_text(text, file_name)
    pdf_document.close()
    summary = response
    if len(response) > max_length and short_response:
            logging.warning("Die Länge der Zusammenfassung könnte zu Problemen führen, daher wird diese in Abschnitten erneut zusammengefasst")
            summary = summarize(response, summary_prompt)
    pdf_stream.close()
    return summary



def describe_pptx(pptx, short_response=False, max_length=3000, store_content=False):
    """
    Describes the content of a PowerPoint presentation file.
    Args:
        pptx (file-like object): The PowerPoint presentation file to be described.
        short_response (bool, optional): If True, the response will be summarized to fit within the max_length. Defaults to False.
        max_length (int, optional): The maximum length of the response. Defaults to 3000.
        store_content (bool, optional): If True, the content of each slide will be saved to a text file. Defaults to False.
    Returns:
        str: A summary of the PowerPoint presentation content.
    Raises:
        Warning: If the presentation has more than 20 slides, a warning is logged about potential token limit issues.
    Notes:
        - The function reads the PowerPoint file, extracts text and image content from each slide, and summarizes it.
        - If the presentation has more than 20 slides, a warning is logged.
        - If store_content is True, the content of each slide is saved to a text file.
        - If the response exceeds max_length and short_response is True, the response is further summarized.
    """
    pptx_bytes = pptx.read()
    print(pptx.name)
    pptx_stream = io.BytesIO(pptx_bytes)
    # Save the uploaded file temporarily
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
        temp_file.write(pptx_stream.read())
        temp_file_path = temp_file.name  # Get the path of the temporary file
    # Load the PowerPoint presentation
    pptx_stream.close()
    prs = Presentation(temp_file_path)
    os.remove(temp_file_path)  # Remove the temporary file
    response = ""
    image_prompt = "Das folgende Bild ist Teil einer Presentation, fasse zentrale Inhalte des Bildes in maximal 100 Worten zusammen, ohne Information hinzuzufügen"
    # Instruction for summarizing the content of a presentation slide in German
    summary_prompt = "du fasst die Inhalte einer presentationsfolie kurz auf deutsch zusammen. Füge keine Informationen hinzu. Vermeide Geplappere und gebe nur die Zusammenfassung zurück. Füge keine Textformatierung hinzu."
    if len(prs.slides) > 20:
        logging.warning("The length of this presentation might cause issues with the token limit.")

    for slide in prs.slides:
        slide_content = ""
        # Position elements
        for shape in slide.shapes:
            # If the shape has text, draw the text
            if hasattr(shape, "text") and shape.text and shape.text != "<#>":
                slide_content += shape.text

            # If the shape is a picture, add the image
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_stream = io.BytesIO(shape.image.blob)  # Read image data from shape
                slide_content += describe_image(image_stream, image_prompt)

        response += summarize(slide_content, summary_prompt)
        if store_content:
            file_name = str(os.path.splitext(pptx.name)[0])+".txt"
            save_text(slide_content, file_name)
    summary = response
    if len(response) > max_length and short_response:
            logging.warning("Die Länge der Zusammenfassung könnte zu Problemen führen, daher wird diese in Abschnitten erneut zusammengefasst")
            summary = summarize(response, summary_prompt)
    return summary



allowed_paths = [".pptx", ".pdf"]


def describe_file(file, short_response=False, max_length=3000, store_content=False):
    """
    Describes the content of a given file based on its suffix.

    Args:
        file (File): The file to be described.
        short_response (bool, optional): If True, returns a shorter description. Defaults to False.
        max_length (int, optional): The maximum length of the description, only applies if short_response. Defaults to 3000.
        store_content (bool, optional): If True, stores the content of the file. Defaults to False.

    Returns:
        str: The description of the file content.

    Raises:
        ValueError: If the file type is not supported or no function is found to describe the file.
    """
    suffix = Path(file.name).suffix
    if suffix in allowed_paths:
        try:
            description_function = globals()[f"describe_{suffix[1:]}"]
            description = description_function(file, short_response, max_length, store_content)
        except KeyError:
            raise ValueError(f"No function found to describe files with suffix {suffix}.")
    else:
        raise ValueError("The file type is not supported.")
    return description

