from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import io
import base64
from groq import Groq
import logging

def encode_image_from_stream(image_stream):
    return base64.b64encode(image_stream.getvalue()).decode('utf-8')



def describe_pptx(pptx_path):
    """
    Summarizes the content of a PowerPoint presentation in German.
    Args:
        pptx_path (str): The file path to the PowerPoint presentation.
    Returns:
        str: A summarized description of the PowerPoint presentation content.
    Raises:
        ValueError: If the GROQ_API_KEY environment variable is not set.
        Exception: If there is an error during the summary API request.
    Notes:
        - The function uses the Groq API to summarize text and image content from each slide.
        - If the presentation has more than 20 slides, a warning is logged due to potential token limit issues.
        - The function handles both text and image content, summarizing images using a specific instruction.
        - If the final summary exceeds 3000 characters, it is further summarized in chunks.
    """
    # Load the PowerPoint presentation
    prs = Presentation(pptx_path)
    slide_content = []
    response = ""
    GROQ_API_KEY = os.getenv('GROQ_API_KEY')
    if GROQ_API_KEY is None:
        raise ValueError("The GROQ_API_KEY environment variable is not set.")
    # Initialize the Groq client
    client = Groq(api_key=GROQ_API_KEY)

    # Instruction for summarizing the content of a presentation slide in German
    summary_instruction = "du fasst die Inhalte einer presentationsfolie kurz auf deutsch zusammen. Füge keine Informationen hinzu. Vermeide Geplappere und gebe nur die Zusammenfassung zurück. Füge keine Textformatierung hinzu."
    if len(prs.slides) > 20:
        logging.warning("The length of this presentation might cause issues with the token limit.")


    for slide_number,slide in enumerate(prs.slides):
        # Position elements
        for shape in slide.shapes:
            # If the shape has text, draw the text
            if hasattr(shape, "text") and shape.text:
                slide_content.append(shape.text)
            
            # If the shape is a picture, add the image
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_stream = io.BytesIO(shape.image.blob)  # Read image data from shape
                base64_image = encode_image_from_stream(image_stream)

                #API request
                chat_completion = client.chat.completions.create(
                    messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text", "text": "Das folgende Bild ist Teil einer Presentation, fasse zentrale Inhalte des Bildes in maximal 100 Worten zusammen, ohne Information hinzuzufügen"
                            },
                            {
                        "type": "image_url",
                        "image_url": {

                            "url": f"data:image/png;base64,{base64_image}",
                        },},
                           ],
                    }
                    ],
                model="llama-3.2-90b-vision-preview",
                    
                )
                slide_content.append(chat_completion.choices[0].message.content)

        #Summarize the descriptions
        summary_input = ""
        
        for content in slide_content:
            if content != '‹#›':
                summary_input = summary_input + " " + content
        try:
          chat_completion = client.chat.completions.create(
            messages=[
                {
                    "role": "system",
                    "content": summary_instruction,
                },
                {
                    "role": "user",
                    "content": summary_input,
                },
            ],
            model="llama3-8b-8192",
          )
          response += chat_completion.choices[0].message.content
        except Exception as e:
          logging.error(f"Error during summary API request: {e}")
        slide_content.clear()

    max_length = 3000
    if len(response) > max_length:
        try:
            logging.warning("Die Länge der Zusammenfassung könnte zu Problemen führen, daher wird diese in Abschnitten erneut zusammengefasst")
            chunks = [response[i:i + max_length] for i in range(0, len(response), max_length)]
            summarized_response = ""
            for chunk in chunks:
                chat_completion = client.chat.completions.create(
                    messages=[
                        {
                            "role": "system",
                            "content": "du fasst die Inhalte einer presentation kurz zusammen. Füge keine Informationen hinzu. Vermeide Geplappere und gebe nur die Zusammenfassung zurück. Füge keine Textformatierung hinzu.",
                        },
                        {
                            "role": "user",
                            "content":  chunk,
                        },
                    ],
                    model="llama3-8b-8192",
                )
                summarized_response += chat_completion.choices[0].message.content
            response = summarized_response
        except Exception as e:
            logging.error(f"Error during summary API request: {e}")

    print(response)
    return response


describe_pptx(r"C:\Users\Jonas\Downloads\schlacht-bei-tannenberg(1).pptx")